from __future__ import annotations

from pathlib import Path
from typing import Any

from hypermindlabs.document_parser.adapters.base import (
    DocumentParseProfile,
    DocumentParserAdapter,
    clamp_confidence,
    clamp_cost,
)

try:
    from docx import Document as DocxDocument
except Exception:  # noqa: BLE001
    DocxDocument = None  # type: ignore[assignment]

try:
    from openpyxl import load_workbook
except Exception:  # noqa: BLE001
    load_workbook = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except Exception:  # noqa: BLE001
    Presentation = None  # type: ignore[assignment]


_OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}


def _runtime_value(config_manager: Any | None, path: str, default: Any) -> Any:
    if config_manager is None:
        return default
    try:
        return config_manager.runtimeValue(path, default)
    except Exception:  # noqa: BLE001
        return default


def _runtime_int(config_manager: Any | None, path: str, default: int) -> int:
    try:
        return int(_runtime_value(config_manager, path, default))
    except (TypeError, ValueError):
        return int(default)


def _safe_text(value: Any) -> str:
    return str(value if value is not None else "").strip()


def _append_section(
    sections: list[dict[str, Any]],
    *,
    text: str,
    element_type: str,
    level: int,
    cursor: int,
    metadata: dict[str, Any] | None = None,
) -> int:
    value = _safe_text(text)
    if not value:
        return cursor
    start_char = int(cursor)
    end_char = start_char + len(value)
    item = {
        "section_id": f"s{len(sections) + 1}",
        "title": value if element_type == "heading" else "",
        "level": max(1, int(level)),
        "text": value,
        "start_char": start_char,
        "end_char": end_char,
        "page_start": None,
        "page_end": None,
        "metadata": {
            "element_type": element_type,
            "source": "office_suite",
        },
    }
    if isinstance(metadata, dict):
        item["metadata"].update(dict(metadata))
    sections.append(item)
    return end_char + 2


class OfficeDocumentParserAdapter(DocumentParserAdapter):
    adapter_name = "office-suite"
    adapter_version = "v1"

    def __init__(self, *, config_manager: Any | None = None):
        self._config = config_manager

    def can_parse(self, profile: DocumentParseProfile) -> bool:
        ext = str(profile.file_extension or "").strip().lower()
        mime = str(profile.file_mime or "").strip().lower()
        if ext in _OFFICE_EXTENSIONS:
            return True
        if mime in {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        }:
            return True
        return False

    def confidence(self, profile: DocumentParseProfile) -> float:
        ext = str(profile.file_extension or "").strip().lower()
        score = 0.2
        if ext == ".docx":
            score += 0.6
        elif ext == ".xlsx":
            score += 0.52
        elif ext == ".pptx":
            score += 0.48
        mime = str(profile.file_mime or "").strip().lower()
        if "openxmlformats" in mime:
            score += 0.2
        return clamp_confidence(score, default=0.0)

    def cost(self, profile: DocumentParseProfile) -> float:
        size_mb = max(0.0, float(profile.file_size_bytes or 0) / (1024.0 * 1024.0))
        return clamp_cost(1.7 + (size_mb / 1.6), default=1.7)

    def _parse_docx(self, path: Path) -> dict[str, Any]:
        if DocxDocument is None:
            raise RuntimeError("python-docx dependency is not available for docx parsing")

        document = DocxDocument(str(path))
        sections: list[dict[str, Any]] = []
        cursor = 0

        for paragraph in document.paragraphs:
            text = _safe_text(paragraph.text)
            if not text:
                continue
            style_name = _safe_text(getattr(paragraph.style, "name", ""))
            lower_style = style_name.lower()
            element_type = "paragraph"
            level = 2
            if lower_style.startswith("heading"):
                element_type = "heading"
                digits = "".join(ch for ch in style_name if ch.isdigit())
                level = int(digits) if digits else 1
            elif "list" in lower_style or text.startswith(("- ", "* ", "• ")):
                element_type = "list"
                level = 2
            cursor = _append_section(
                sections,
                text=text,
                element_type=element_type,
                level=level,
                cursor=cursor,
                metadata={"style": style_name},
            )

        for table_index, table in enumerate(document.tables, start=1):
            rows: list[str] = []
            for row in table.rows:
                cells = [_safe_text(cell.text) for cell in row.cells]
                cells = [cell for cell in cells if cell]
                if cells:
                    rows.append(" | ".join(cells))
            table_text = "\n".join(rows).strip()
            cursor = _append_section(
                sections,
                text=table_text,
                element_type="table",
                level=2,
                cursor=cursor,
                metadata={"table_index": table_index, "row_count": len(rows)},
            )

        content_text = "\n\n".join(section.get("text", "") for section in sections if section.get("text"))
        return {
            "status": "parsed" if content_text else "partial",
            "content_text": content_text,
            "sections": sections,
            "metadata": {
                "parser_kind": "office_suite",
                "office_format": "docx",
                "paragraph_count": len(document.paragraphs),
                "table_count": len(document.tables),
            },
            "warnings": ["office_docx_empty_content"] if not content_text else [],
            "errors": [],
        }

    def _parse_xlsx(self, path: Path) -> dict[str, Any]:
        if load_workbook is None:
            raise RuntimeError("openpyxl dependency is not available for xlsx parsing")

        max_rows = max(1, _runtime_int(self._config, "documents.office_xlsx_max_rows", 300))
        max_cols = max(1, _runtime_int(self._config, "documents.office_xlsx_max_cols", 20))

        workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
        sections: list[dict[str, Any]] = []
        cursor = 0
        rows_parsed = 0
        sheets_parsed = 0
        truncated_rows = False

        for sheet in workbook.worksheets:
            sheets_parsed += 1
            cursor = _append_section(
                sections,
                text=f"Sheet: {sheet.title}",
                element_type="heading",
                level=2,
                cursor=cursor,
                metadata={"sheet": sheet.title},
            )
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > max_rows:
                    truncated_rows = True
                    break
                values = [_safe_text(cell) for cell in list(row)[:max_cols]]
                if not any(values):
                    continue
                rows_parsed += 1
                row_text = " | ".join(value for value in values if value)
                cursor = _append_section(
                    sections,
                    text=row_text,
                    element_type="table",
                    level=2,
                    cursor=cursor,
                    metadata={
                        "sheet": sheet.title,
                        "row_index": row_index,
                        "column_count": len(values),
                    },
                )

        workbook.close()

        content_text = "\n\n".join(section.get("text", "") for section in sections if section.get("text"))
        warnings: list[str] = []
        if truncated_rows:
            warnings.append("office_xlsx_row_limit_applied")
        if not content_text:
            warnings.append("office_xlsx_empty_content")

        return {
            "status": "parsed" if content_text else "partial",
            "content_text": content_text,
            "sections": sections,
            "metadata": {
                "parser_kind": "office_suite",
                "office_format": "xlsx",
                "sheets_parsed": sheets_parsed,
                "rows_parsed": rows_parsed,
                "max_rows": max_rows,
                "max_cols": max_cols,
            },
            "warnings": warnings,
            "errors": [],
        }

    def _parse_pptx(self, path: Path) -> dict[str, Any]:
        if Presentation is None:
            raise RuntimeError("python-pptx dependency is not available for pptx parsing")

        max_slides = max(1, _runtime_int(self._config, "documents.office_pptx_max_slides", 200))
        presentation = Presentation(str(path))
        sections: list[dict[str, Any]] = []
        cursor = 0
        slide_count = 0
        truncated = False

        for slide_index, slide in enumerate(presentation.slides, start=1):
            if slide_index > max_slides:
                truncated = True
                break
            slide_count += 1
            cursor = _append_section(
                sections,
                text=f"Slide {slide_index}",
                element_type="heading",
                level=2,
                cursor=cursor,
                metadata={"slide_index": slide_index},
            )
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = _safe_text(getattr(shape, "text", ""))
                if not text:
                    continue
                cursor = _append_section(
                    sections,
                    text=text,
                    element_type="paragraph",
                    level=2,
                    cursor=cursor,
                    metadata={"slide_index": slide_index, "shape_index": shape_index},
                )

        content_text = "\n\n".join(section.get("text", "") for section in sections if section.get("text"))
        warnings: list[str] = []
        if truncated:
            warnings.append("office_pptx_slide_limit_applied")
        if not content_text:
            warnings.append("office_pptx_empty_content")

        return {
            "status": "parsed" if content_text else "partial",
            "content_text": content_text,
            "sections": sections,
            "metadata": {
                "parser_kind": "office_suite",
                "office_format": "pptx",
                "slides_parsed": slide_count,
                "max_slides": max_slides,
            },
            "warnings": warnings,
            "errors": [],
        }

    def parse(self, profile: DocumentParseProfile) -> dict[str, Any]:
        path = Path(profile.file_path)
        if not path.exists():
            raise FileNotFoundError(f"Parser input file does not exist: {path}")

        ext = str(profile.file_extension or "").strip().lower()
        try:
            if ext == ".docx":
                output = self._parse_docx(path)
            elif ext == ".xlsx":
                output = self._parse_xlsx(path)
            elif ext == ".pptx":
                output = self._parse_pptx(path)
            else:
                return {
                    "status": "failed",
                    "content_text": "",
                    "sections": [],
                    "metadata": {
                        "parser_kind": "office_suite",
                        "office_format": ext,
                    },
                    "warnings": [],
                    "errors": [f"unsupported_office_format:{ext or 'unknown'}"],
                    "confidence": 0.0,
                    "cost": self.cost(profile),
                }
        except Exception as error:  # noqa: BLE001
            return {
                "status": "failed",
                "content_text": "",
                "sections": [],
                "metadata": {
                    "parser_kind": "office_suite",
                    "office_format": ext,
                },
                "warnings": [],
                "errors": [f"office_parse_error:{error}"],
                "confidence": 0.0,
                "cost": self.cost(profile),
            }

        output["confidence"] = self.confidence(profile)
        output["cost"] = self.cost(profile)
        return output
